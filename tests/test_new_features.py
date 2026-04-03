"""Tests for the new features: PPTX extraction, per-model AI detection, word quotas."""

from __future__ import annotations

import io
from unittest.mock import AsyncMock, patch

import pytest

from app.tools.content_extractor_tool import extract_text, _extract_from_pptx
from app.tools.ai_detection_tool import detect_ai_text
from app.services.rate_limiter import limiter, UserTier
from app.models.schemas import FileType


# ---------------------------------------------------------------------------
# 1. PPTX Extraction
# ---------------------------------------------------------------------------

class TestPPTXExtraction:
    def _make_pptx_bytes(self, texts: list[str]) -> bytes:
        """Create a minimal PPTX file in memory with the given slide texts."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        for text in texts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
            txBox.text_frame.text = text
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    @pytest.mark.asyncio
    async def test_extract_pptx_basic(self):
        """Should extract text from a PPTX file."""
        pptx_bytes = self._make_pptx_bytes(["Hello World", "Second Slide"])
        result = await extract_text(pptx_bytes, "test.pptx")
        assert "Hello World" in result["text"]
        assert "Second Slide" in result["text"]
        assert result["file_type"] == "pptx"

    @pytest.mark.asyncio
    async def test_extract_pptx_empty(self):
        """Empty PPTX should raise ValueError."""
        pptx_bytes = self._make_pptx_bytes(["", ""])
        with pytest.raises(ValueError, match="empty"):
            await extract_text(pptx_bytes, "empty.pptx")

    def test_filetype_enum_has_pptx(self):
        """FileType enum should include PPTX."""
        assert FileType.PPTX.value == "pptx"

    def test_pptx_in_allowed_extensions(self):
        from app.config import settings
        assert ".pptx" in settings.allowed_extensions


# ---------------------------------------------------------------------------
# 2. Per-Model AI Detection
# ---------------------------------------------------------------------------

class TestPerModelDetection:
    @pytest.mark.asyncio
    async def test_model_attribution_in_gpt_results(self):
        """GPT classification should include model_attribution in indicators."""
        chunks = ["AI-generated text that sounds very robotic. " * 5] * 3
        text = " ".join(chunks)

        gpt_results = [
            {"label": "ai", "confidence": 0.9, "reason": "Formulaic", "suspected_model": "chatgpt"},
            {"label": "ai", "confidence": 0.85, "reason": "Hedging", "suspected_model": "claude"},
            {"label": "human", "confidence": 0.8, "reason": "Natural", "suspected_model": "human"},
        ]

        with patch("app.tools.ai_detection_tool._gpt_classify_chunks", new_callable=AsyncMock) as mock_gpt:
            mock_gpt.return_value = gpt_results
            result = await detect_ai_text(text, chunks=chunks, use_gpt=True)

        attr = result["indicators"]["model_attribution"]
        assert attr is not None
        assert attr.get("chatgpt") == 1
        assert attr.get("claude") == 1
        assert "human" not in attr  # only AI labels counted

    @pytest.mark.asyncio
    async def test_model_attribution_none_without_gpt(self):
        """Without GPT mode, model_attribution should be None."""
        text = "Normal text for heuristic analysis. " * 10
        result = await detect_ai_text(text)
        assert result["indicators"]["model_attribution"] is None

    @pytest.mark.asyncio
    async def test_flagged_chunks_include_suspected_model(self):
        """Flagged chunks from GPT should include suspected_model."""
        chunks = ["Some AI text here. " * 10]
        text = chunks[0]

        gpt_results = [
            {"label": "ai", "confidence": 0.9, "reason": "Robotic", "suspected_model": "gemini"},
        ]

        with patch("app.tools.ai_detection_tool._gpt_classify_chunks", new_callable=AsyncMock) as mock_gpt:
            mock_gpt.return_value = gpt_results
            result = await detect_ai_text(text, chunks=chunks, use_gpt=True)

        gpt_flagged = [f for f in result["flagged_chunks"] if "GPT classifier" in f.get("reason", "")]
        assert len(gpt_flagged) >= 1
        assert gpt_flagged[0]["suspected_model"] == "gemini"


# ---------------------------------------------------------------------------
# 3. Word-Count Quotas
# ---------------------------------------------------------------------------

class TestWordQuotas:
    def test_word_quota_config(self):
        from app.config import settings
        assert settings.word_quota_free == 5000
        assert settings.word_quota_pro == 200000
        assert settings.word_quota_premium == 500000

    def test_word_quota_for_tier(self):
        from app.config import settings
        assert limiter._word_quota_for_tier(UserTier.FREE) == settings.word_quota_free
        assert limiter._word_quota_for_tier(UserTier.PRO) == settings.word_quota_pro
        assert limiter._word_quota_for_tier(UserTier.PREMIUM) == settings.word_quota_premium

    def test_check_word_quota_allowed(self):
        """Fresh user should have full quota available."""
        with patch.object(limiter, "get_monthly_word_count", return_value=0):
            result = limiter.check_word_quota(1, UserTier.FREE, word_count=100)
            assert result["allowed"] is True
            assert result["remaining"] == 5000

    def test_check_word_quota_exceeded(self):
        """User over quota should be denied."""
        with patch.object(limiter, "get_monthly_word_count", return_value=5000):
            result = limiter.check_word_quota(1, UserTier.FREE, word_count=100)
            assert result["allowed"] is False
            assert result["remaining"] == 0

    def test_record_usage_accepts_word_count(self):
        """record_usage should accept word_count parameter without error."""
        with patch("app.services.rate_limiter.UsageRateLimiter.get_count", return_value=0):
            with patch("app.services.database.get_db") as mock_db:
                mock_db.return_value.execute.return_value = None
                # Should not raise
                limiter.record_usage(
                    "user:1", "scan",
                    user_id=1, ip_address="127.0.0.1", word_count=250,
                )
                # Verify word_count was passed to the SQL
                call_args = mock_db.return_value.execute.call_args
                assert 250 in call_args[0][1]  # word_count in params

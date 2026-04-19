"""Content extractor tool — extracts text from PDF, DOCX, TXT, and LaTeX files.

Standalone, framework-agnostic tool. Returns structured JSON.
Also provides text chunking functionality.
"""

from __future__ import annotations

import re
import time
from collections.abc import Callable
from io import BytesIO
from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from PyPDF2 import PdfReader

from app.config import settings
from app.utils.logger import get_logger

logger = get_logger(__name__)


async def extract_text(file_bytes: bytes, filename: str) -> dict:
    """Extract text content from an uploaded file.

    Args:
        file_bytes: Raw bytes of the uploaded file.
        filename: Original filename (used to determine type).

    Returns:
        Dict with ``filename``, ``file_type``, ``text``, ``char_count``, ``elapsed_s``.

    Raises:
        ValueError: If the file type is unsupported or extraction fails.
    """
    max_bytes = getattr(settings, "max_upload_size_mb", 50) * 1024 * 1024
    if len(file_bytes) > max_bytes:
        raise ValueError(
            f"File too large: {len(file_bytes)} bytes exceeds limit of {max_bytes}"
        )

    ext = Path(filename).suffix.lower()
    start = time.perf_counter()

    extractors: dict[str, Callable[[bytes], str]] = {
        ".pdf": _extract_from_pdf,
        ".docx": _extract_from_docx,
        ".txt": _extract_from_txt,
        ".tex": _extract_from_latex,
        ".pptx": _extract_from_pptx,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {ext}")

    logger.info("extracting_text", filename=filename, file_type=ext)
    import asyncio
    loop = asyncio.get_running_loop()
    text, fallback_used = await loop.run_in_executor(None, extractor, file_bytes)

    if not text.strip():
        raise ValueError("Extracted text is empty — the file may be scanned or corrupt.")

    elapsed = round(time.perf_counter() - start, 3)

    warnings_list: list[str] = []
    if fallback_used:
        warnings_list.append("fallback_extractor_used")

    quality = "high"
    if fallback_used:
        quality = "medium"

    logger.info(
        "extraction_complete",
        filename=filename,
        char_count=len(text),
        elapsed_s=elapsed,
    )

    return {
        "filename": filename,
        "file_type": ext.lstrip("."),
        "text": text,
        "char_count": len(text),
        "elapsed_s": elapsed,
        "extraction_quality": quality,
        "fallback_used": fallback_used,
        "warnings": warnings_list,
    }


def chunk_text(
    text: str,
    chunk_size: int | None = None,
    overlap: int | None = None,
) -> dict:
    """Split text into overlapping chunks.

    Uses sentence-boundary-aware splitting.

    Args:
        text: The full document text.
        chunk_size: Target chars per chunk (default from settings).
        overlap: Overlap chars between chunks (default from settings).

    Returns:
        Dict with ``chunks``, ``chunk_count``, ``chunk_size``, ``overlap``.
    """
    if chunk_size is None:
        chunk_size = settings.chunk_size
    if overlap is None:
        overlap = settings.chunk_overlap

    if overlap >= chunk_size:
        raise ValueError(f"overlap ({overlap}) must be less than chunk_size ({chunk_size})")

    start = time.perf_counter()
    text = text.strip()

    if not text:
        return {"chunks": [], "chunk_count": 0, "chunk_size": chunk_size, "overlap": overlap, "elapsed_s": 0.0}

    if len(text) <= chunk_size:
        return {"chunks": [text], "chunk_count": 1, "chunk_size": chunk_size, "overlap": overlap, "elapsed_s": 0.0}

    chunks: list[str] = []
    pos = 0

    while pos < len(text):
        end = pos + chunk_size

        if end < len(text):
            boundary = _find_sentence_boundary(text, pos, end)
            if boundary > pos:
                end = boundary

        chunk = text[pos:end].strip()
        if chunk:
            chunks.append(chunk)

        step = max(end - pos - overlap, 1)
        pos += step

        # Align to word boundary — never start a chunk mid-word
        if 0 < pos < len(text) and not text[pos - 1].isspace():
            scan_start = pos
            while pos < len(text) and not text[pos].isspace():
                pos += 1
            # Skip past the whitespace to the first char of the next word
            while pos < len(text) and text[pos].isspace():
                pos += 1
            # Safety: revert if we had to skip more than 50 chars
            if pos - scan_start > 50:
                pos = end

    elapsed = round(time.perf_counter() - start, 3)

    logger.info(
        "text_chunked",
        total_chunks=len(chunks),
        chunk_size=chunk_size,
        overlap=overlap,
        elapsed_s=elapsed,
    )

    return {
        "chunks": chunks,
        "chunk_count": len(chunks),
        "chunk_size": chunk_size,
        "overlap": overlap,
        "elapsed_s": elapsed,
    }


# ---------------------------------------------------------------------------
# Private helpers
# ---------------------------------------------------------------------------

def _extract_from_pdf(file_bytes: bytes) -> tuple[str, bool]:
    # Primary: pdfplumber (much better at preserving spaces)
    try:
        import pdfplumber
        pages: list[str] = []
        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    pages.append(page_text)
        text = "\n".join(pages)
        if text.strip():
            return _fix_pdf_spacing(text), False
    except Exception as exc:
        logger.warning("pdfplumber_failed_falling_back_to_pypdf2", error=str(exc))

    # Fallback: PyPDF2
    reader = PdfReader(BytesIO(file_bytes))
    pages = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            pages.append(page_text)
    return _fix_pdf_spacing("\n".join(pages)), True


# Short function words that commonly get stuck to a preceding word in PDFs.
# Only 2-3 char words to keep false-positive risk near zero.
_SUFFIX_WORDS = frozenset([
    "in", "on", "of", "at", "by", "an", "as", "or", "to", "is", "it",
    "and", "the", "for", "are", "but", "not", "was", "has", "its",
])

# Match word tokens of 5+ chars (shorter tokens are unlikely merges)
_SUFFIX_JOIN_RE = re.compile(r'\b([A-Za-z]{5,})\b')


def _split_suffix_joins(text: str) -> str:
    """Split tokens where a known word has a short function word glued to the end.

    E.g. "Layersin" → "Layers in", "Organicand" → "Organic and".
    Only splits when the prefix (minus suffix) is in the word set.
    """
    def _try(m: re.Match) -> str:
        token = m.group(1)
        lower = token.lower()
        # Don't touch tokens that are already known words
        if lower in _WORD_SET:
            return token
        # Try stripping 2-char then 3-char suffixes
        for slen in (2, 3):
            if len(lower) <= slen + 2:
                continue
            suffix = lower[-slen:]
            prefix = lower[:-slen]
            if suffix in _SUFFIX_WORDS and prefix in _WORD_SET:
                return token[:-slen] + " " + suffix
        return token

    return _SUFFIX_JOIN_RE.sub(_try, text)


def _fix_pdf_spacing(text: str) -> str:
    """Post-process extracted PDF text to restore missing spaces.

    Uses a two-stage approach: safe fixes always run, expensive DP-based
    splitting only runs when the text shows signs of severe merge artifacts.
    """
    # --- Stage 1: safe fixes (always apply) ---
    # Insert space between a letter and an opening paren
    text = re.sub(r'([a-zA-Z])\(', r'\1 (', text)
    # Insert space between a closing paren/period/comma and a letter OR digit
    text = re.sub(r'([.,;:)])([a-zA-Z0-9])', r'\1 \2', text)
    # Insert space between a lowercase letter and an uppercase letter
    # (camelCase joins like "MetalOxideNanocrystals" → "Metal Oxide Nanocrystals").
    # This is safe even for clean text — legitimate camelCase tokens in
    # academic PDFs (e.g. "iPhone") are extremely rare.
    text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
    # Split tokens where a known word has a short function word stuck to the end
    # (e.g. "Layersin" → "Layers in", "Organicand" → "Organic and").
    # Only fires when prefix is a known word and suffix is a very common
    # short function word — safe because the intersection is tiny.
    text = _split_suffix_joins(text)
    # Insert space between a digit and a letter (but not inside known patterns like "3D")
    text = re.sub(r'(\d)([a-zA-Z]{2,})', r'\1 \2', text)
    # Insert space between a letter and a 4-digit year (e.g. "Interfaces2021")
    text = re.sub(r'([a-zA-Z])(\d{4})\b', r'\1 \2', text)
    # Collapse multiple spaces
    text = re.sub(r' {2,}', ' ', text)

    # --- Stage 2: DP-based splitting for remaining lowercase merges ---
    # Always run — the _MERGE_RE pattern (12+ lowercase chars) is specific
    # enough that it rarely matches real words, and the 200-blob cap
    # prevents CPU spikes.
    text = _split_merged_words(text)

    # --- Stage 3: expensive capitalized-join splitting (only when severely broken) ---
    word_count = max(text.count(' '), 1)
    non_space_chars = len(text.replace(' ', ''))
    avg_word_len = non_space_chars / word_count
    if avg_word_len > 8:
        # Split tokens where a capitalized prefix is joined to a known word
        text = _split_capitalized_joins(text)
        # Collapse multiple spaces again
        text = re.sub(r' {2,}', ' ', text)
    return text


# Pattern: a capitalized word token of 6+ chars — short words like "Under"
# are almost never merge artifacts so skip them for performance.
_CAP_JOIN_RE = re.compile(r'\b([A-Z][a-z]{5,})\b')


def _split_capitalized_joins(text: str) -> str:
    """Split tokens like 'Inthis' -> 'In this', 'Interfaciallayers' -> 'Interfacial layers'."""

    def _try_cap_split(m: re.Match) -> str:
        word = m.group(1)
        lower = word.lower()

        # Check if the full word itself is known (don't split real words)
        if lower in _WORD_SET:
            return m.group(0)

        # Only try splitting if the word is plausibly two+ words merged
        # (at least one valid prefix of 2+ chars that is a known word)
        has_known_prefix = False
        for i in range(2, min(len(lower), _MAX_WORD_LEN + 1)):
            if lower[:i] in _WORD_SET:
                has_known_prefix = True
                break
        if not has_known_prefix:
            return m.group(0)

        # Try all split points: prefix (capitalized) + suffix (lowercase)
        for i in range(2, len(lower) - 1):
            prefix = lower[:i]
            suffix = lower[i:]
            if prefix in _WORD_SET and suffix in _WORD_SET:
                return word[:i] + " " + suffix

            if prefix in _WORD_SET:
                split_suffix = _dp_split(suffix)
                if split_suffix:
                    return word[:i] + " " + split_suffix

        # Try DP on the full lowercase blob
        split_full = _dp_split(lower)
        if split_full:
            words = split_full.split()
            words[0] = words[0].capitalize()
            return " ".join(words)

        return m.group(0)

    return _CAP_JOIN_RE.sub(_try_cap_split, text)


# ---------------------------------------------------------------------------
# Word list for merge-point detection (Issue 7: load from file, fallback to inline)
# ---------------------------------------------------------------------------
_WORD_LIST_PATH = Path(__file__).resolve().parent / "common_words.txt"


def _load_word_set() -> frozenset[str]:
    """Load words from common_words.txt if available, else use inline fallback."""
    if _WORD_LIST_PATH.is_file():
        try:
            words = _WORD_LIST_PATH.read_text(encoding="utf-8").splitlines()
            words = [w.strip().lower() for w in words if w.strip() and not w.startswith("#")]
            if len(words) > 100:
                return frozenset(words)
        except Exception:
            pass

    # Inline fallback — general English + academic terms (no duplicates)
    return frozenset([
        "of", "in", "on", "at", "by", "an", "as", "or", "so", "if", "to",
        "no", "do", "up", "we", "he",
        "the", "and", "for", "are", "but", "not", "you", "all", "can", "had",
        "her", "was", "one", "our", "out", "has", "his", "how", "its", "may",
        "new", "now", "old", "see", "way", "who", "did", "get", "let", "say",
        "she", "too", "use",
        "with", "have", "from", "this", "that", "been", "into", "also", "both",
        "some", "than", "them", "then", "were", "when", "will", "each", "made",
        "more", "most", "much", "must", "only", "over", "such", "take", "very",
        "back", "even", "give", "just", "like", "long", "make", "many", "well",
        "what", "your", "used",
        "about", "after", "being", "could", "every", "first", "found", "great",
        "other", "their", "there", "these", "think", "those", "three", "under",
        "water", "where", "which", "while", "world", "would",
        "should", "through", "between", "because", "before", "during",
        "without", "including", "equally", "important", "however", "therefore",
        "resulting", "recently", "relatively", "significantly", "respectively",
        "respect", "order", "effect", "effects", "report", "show", "shown",
        "various", "different", "compared", "observed", "obtained", "proposed",
        "present", "increase", "decrease", "improved", "lower", "higher",
        "using", "work", "paper", "figure", "table", "note", "type",
        "interface", "engineering", "active", "layer", "morphology",
        "organic", "solar", "cells", "cell", "performance", "materials",
        "material", "optimization", "electron", "transport", "transporting",
        "device", "devices", "efficiency", "recent", "advances", "inorganic",
        "hole", "extraction", "bulk", "heterojunction", "photoactive",
        "respective", "electrodes", "electrode", "review", "summarize",
        "progress", "based", "high", "power", "conversion", "exceeding",
        "resulted", "advantages", "synthetic", "versatility", "absorption",
        "coefficient", "wavelength", "range", "thermal", "stability",
        "attained", "interfacial", "layers", "emergence", "nonfullerene",
        "small", "molecule", "acceptors", "acceptor", "facilitating",
        "concept", "basic", "application", "structure", "properties",
        "analysis", "results", "system", "study", "method", "model",
        "data", "process", "energy", "surface", "film", "optical",
        "characterization", "measurement", "fabrication", "deposition",
        "temperature", "annealing", "photovoltaic", "cathode", "anode",
        "polymer", "blend", "donor", "charge", "transfer",
        "recombination", "mobility", "current", "density", "voltage",
        "spectrum", "response", "quantum", "coupling", "interaction",
    ])


_WORD_SET = _load_word_set()
_MAX_WORD_LEN = max((len(w) for w in _WORD_SET), default=0)

# Match whole tokens of 14+ alphabetic characters — likely merge artifacts.
# Works on word boundaries so we get the full token for DP splitting.
_MERGE_RE = re.compile(r'\b([a-zA-Z]{14,})\b')

# Cap the number of blobs processed per document to avoid CPU spikes.
_MAX_BLOB_SPLITS = 200


def _split_merged_words(text: str) -> str:
    """Find long tokens (14+ chars) and try to split them into words."""
    count = 0

    def _try_split(match: re.Match) -> str:
        nonlocal count
        if count >= _MAX_BLOB_SPLITS:
            return match.group()
        count += 1
        token = match.group(1)
        lower = token.lower()
        # Skip tokens that are already known words
        if lower in _WORD_SET:
            return match.group()
        result = _dp_split(lower)
        if result:
            # Preserve leading capitalisation from the original token
            if token[0].isupper():
                words = result.split()
                words[0] = words[0].capitalize()
                return " ".join(words)
            return result
        return match.group()

    return _MERGE_RE.sub(_try_split, text)


def _dp_split(blob: str) -> str | None:
    """Use DP to find the segmentation that maximizes known-word coverage.

    Returns space-separated words if >=60% of chars are covered.
    Does NOT break on the first matching word — explores all lengths so that
    the globally best split is found (Issue 6).
    """
    n = len(blob)
    if n == 0:
        return None

    NEG_INF = float("-inf")

    # dp[i] = max covered chars for blob[0:i], parent for traceback
    dp = [NEG_INF] * (n + 1)
    dp[0] = 0
    parent: list[tuple[int, bool]] = [(-1, False)] * (n + 1)

    for i in range(n):
        if dp[i] == NEG_INF:
            continue

        # Try ALL known words starting at i (no break — pick the best globally)
        for wlen in range(min(_MAX_WORD_LEN, n - i), 0, -1):
            candidate = blob[i:i + wlen]
            if candidate in _WORD_SET:
                new_cov = dp[i] + wlen
                if new_cov > dp[i + wlen]:
                    dp[i + wlen] = new_cov
                    parent[i + wlen] = (i, True)

        # Also allow consuming one char as unknown
        if dp[i] > dp[i + 1]:
            dp[i + 1] = dp[i]
            parent[i + 1] = (i, False)

    if dp[n] == NEG_INF or dp[n] / n < 0.6:
        return None

    # Traceback
    segments: list[str] = []
    pos = n
    while pos > 0:
        start, _ = parent[pos]
        segments.append(blob[start:pos])
        pos = start
    segments.reverse()

    # Merge adjacent unknown fragments
    merged: list[str] = []
    for seg in segments:
        if seg not in _WORD_SET and merged and merged[-1] not in _WORD_SET:
            merged[-1] += seg
        else:
            merged.append(seg)

    if len(merged) <= 1:
        return None

    # Validate the split quality:
    # 1. If ALL segments are known words and there are 2+, accept (perfect split)
    # 2. Otherwise require at least 2 substantive (≥3 chars) known words
    #    to avoid false splits like "infrastructure" → "in fr a structure"
    known_count = sum(1 for s in merged if s in _WORD_SET)
    all_known = all(s in _WORD_SET for s in merged)
    if all_known and known_count >= 2:
        pass  # perfect split — accept
    else:
        substantive = sum(1 for s in merged if s in _WORD_SET and len(s) >= 3)
        if substantive < 2:
            return None
        # If unknown fragments are longer than 2 chars each on average, bail
        unknown_segs = [s for s in merged if s not in _WORD_SET]
        unknown_chars = sum(len(s) for s in unknown_segs)
        if unknown_segs and unknown_chars / len(unknown_segs) > 2:
            return None

    return " ".join(merged)


def _extract_from_docx(file_bytes: bytes) -> tuple[str, bool]:
    doc = DocxDocument(BytesIO(file_bytes))
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip()), False


def _extract_from_txt(file_bytes: bytes) -> tuple[str, bool]:
    return file_bytes.decode("utf-8", errors="replace"), False


def _extract_from_pptx(file_bytes: bytes) -> tuple[str, bool]:
    """Extract text from a PowerPoint (.pptx) file."""
    prs = Presentation(BytesIO(file_bytes))
    slides_text: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            parts.append(text)
        if parts:
            slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text), False


def _extract_from_latex(file_bytes: bytes) -> tuple[str, bool]:
    """Extract plain text from a LaTeX (.tex) file by stripping commands."""
    raw = file_bytes.decode("utf-8", errors="replace")

    # Remove comments (lines starting with %)
    raw = re.sub(r"(?m)%.*$", "", raw)

    # Remove common preamble commands
    raw = re.sub(r"\\documentclass(\[.*?\])?\{.*?\}", "", raw)
    raw = re.sub(r"\\usepackage(\[.*?\])?\{.*?\}", "", raw)
    raw = re.sub(r"\\(begin|end)\{(document|abstract|figure|table|equation|align|itemize|enumerate|description|thebibliography|verbatim|lstlisting|tabular|array)\}", "", raw)

    # Remove label, ref, cite kept as text: [AuthorYear] style
    raw = re.sub(r"\\label\{[^}]*\}", "", raw)
    raw = re.sub(r"\\(ref|eqref|pageref)\{[^}]*\}", "[ref]", raw)
    raw = re.sub(r"\\(cite|citep|citet|parencite|textcite|autocite)(\[[^\]]*\])?\{([^}]*)\}", r"[\3]", raw)

    # Preserve text from formatting commands (loop handles nested braces)
    _fmt_re = re.compile(r"\\(textbf|textit|emph|underline|texttt|textrm|textsf|textsc)\{([^}]*)\}")
    prev = None
    while prev != raw:
        prev = raw
        raw = _fmt_re.sub(r"\2", raw)
    _section_re = re.compile(r"\\(title|author|date|chapter|section|subsection|subsubsection|paragraph|subparagraph)\*?\{([^}]*)\}")
    prev = None
    while prev != raw:
        prev = raw
        raw = _section_re.sub(r"\2", raw)
    raw = re.sub(r"\\(footnote|footnotetext)\{([^}]*)\}", r" \2", raw)
    raw = re.sub(r"\\caption\{([^}]*)\}", r"\1", raw)

    # Remove math environments but keep simple inline math text
    raw = re.sub(r"\$\$.*?\$\$", " [equation] ", raw, flags=re.DOTALL)
    raw = re.sub(r"\$([^$]+)\$", r"\1", raw)
    raw = re.sub(r"\\begin\{(equation|align|gather|multline)\*?\}.*?\\end\{\1\*?\}", " [equation] ", raw, flags=re.DOTALL)

    # Remove remaining LaTeX commands but keep their text arguments
    raw = re.sub(r"\\[a-zA-Z]+\*?(\{[^}]*\})?", "", raw)

    # Clean up braces and special chars
    raw = raw.replace("{", "").replace("}", "")
    raw = raw.replace("~", " ").replace("``", '"').replace("''", '"')
    raw = re.sub(r"\\[&%$#_]", "", raw)

    # Normalise whitespace
    raw = re.sub(r"[ \t]+", " ", raw)
    raw = re.sub(r"\n{3,}", "\n\n", raw)

    return raw.strip(), False


def _find_sentence_boundary(text: str, start: int, end: int) -> int:
    """Find position just after the last sentence-ending char in [start, end)."""
    best = start
    for sep in (".\n", "\n\n", ". ", ".\t"):
        pos = text.rfind(sep, start, end)
        if pos != -1:
            candidate = pos + len(sep)
            if candidate > best:
                best = candidate
    if best == start:
        pos = text.rfind(".", start, end)
        if pos != -1:
            best = pos + 1
    return best
